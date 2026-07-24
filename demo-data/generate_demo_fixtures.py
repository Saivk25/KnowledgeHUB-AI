"""
Generates the Milestone 12 multi-format seed fixtures, extending the
Milestone 3 convention already established by generate_demo_pdfs.py: a
one-off content generator (not run by CI or any test) that produces
checked-in files under demo-data/, run locally with
`python3 generate_demo_fixtures.py`.

Requires python-docx, python-pptx, and Pillow (all already backend
dependencies -- see apps/api/requirements.txt).

Design note (MILESTONE_12.md Section 4.3): all five files below share the
filename stem `data_retention_policy`, deliberately -- LocalConceptLinker
(app/services/concept_linking.py), the zero-config default, names a new
concept from a resource's filename whenever no `subject` was classified
(the fallback path every LocalHeuristicClassifier-classified resource in
this demo corpus takes), via `fallback_concept_name_from_filename()`:
strip the extension, replace `_`/`-` with spaces, title-case. Every file
here therefore proposes the identical candidate name, "Data Retention
Policy" -- app/services/concept_graph.py's dedup step (exact
normalized-name match) resolves all five to one Concept, giving that one
concept real evidence from five different source types (DOCX, PPTX,
Markdown, Python, PNG/OCR) once ingested through demo-data/seed.py. This
is the same mechanism test_concept_linking_ingestion.py's
`test_two_uploads_with_equivalent_filenames_dedup_to_one_concept` already
proves for two PDFs -- applied here deliberately, across formats, as this
milestone's cross-format concept-linking demonstration.

The YouTube source type is intentionally not generated here -- there is
no file to check into version control for a hosted video, and fetching a
real transcript is a live network call, not a deterministic build step.
See demo-data/YOUTUBE_REFERENCE.md and demo-data/seed.py for how that
source type is represented instead.
"""

from pathlib import Path

import docx
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

OUT_DIR = Path(__file__).parent

# The same retention facts, deliberately repeated (in each format's own
# natural voice) across every fixture below -- this is what makes the
# resulting concept's evidence genuinely cohesive rather than five
# unrelated files that merely happen to share a filename.
RETENTION_FACTS = {
    "financial": ("Financial records", 7, "years", "tax and audit compliance"),
    "employee": ("Employee records", 3, "years after termination", "employment-dispute liability windows"),
    "support": ("Customer support tickets", 2, "years", "product-quality trend analysis"),
}


def _fact_sentence(key: str) -> str:
    label, duration, unit, reason = RETENTION_FACTS[key]
    return f"{label} are retained for {duration} {unit}, primarily to support {reason}."


# ---------------------------------------------------------------------------
# DOCX -- a formal policy document, matching generate_demo_pdfs.py's existing
# Expense_Policy.pdf in register and structure.
# ---------------------------------------------------------------------------


def build_docx() -> None:
    document = docx.Document()
    document.add_heading("Data Retention Policy", level=1)
    document.add_paragraph(
        "This policy defines how long the company retains different categories "
        "of business records before they are securely deleted."
    )
    document.add_heading("Retention Schedule", level=2)
    for key in ("financial", "employee", "support"):
        document.add_paragraph(_fact_sentence(key))
    document.add_heading("Deletion Process", level=2)
    document.add_paragraph(
        "Records past their retention window are purged by an automated nightly "
        "job rather than a manual review, so the schedule above is enforced "
        "consistently regardless of record volume."
    )
    document.save(str(OUT_DIR / "data_retention_policy.docx"))


# ---------------------------------------------------------------------------
# PPTX -- a short slide deck summarizing the same policy.
# ---------------------------------------------------------------------------


def build_pptx() -> None:
    presentation = Presentation()
    title_layout = presentation.slide_layouts[0]
    bullet_layout = presentation.slide_layouts[1]

    title_slide = presentation.slides.add_slide(title_layout)
    title_slide.shapes.title.text = "Data Retention Policy"
    title_slide.placeholders[1].text = "Overview for new hires"

    why_slide = presentation.slides.add_slide(bullet_layout)
    why_slide.shapes.title.text = "Why Retention Matters"
    why_body = why_slide.placeholders[1].text_frame
    why_body.text = "Compliance obligations"
    why_body.add_paragraph().text = "Storage cost control"
    why_body.add_paragraph().text = "Reduced exposure if data is ever breached"

    schedule_slide = presentation.slides.add_slide(bullet_layout)
    schedule_slide.shapes.title.text = "Retention Schedule"
    schedule_body = schedule_slide.placeholders[1].text_frame
    schedule_body.text = _fact_sentence("financial")
    schedule_body.add_paragraph().text = _fact_sentence("employee")
    schedule_body.add_paragraph().text = _fact_sentence("support")

    deletion_slide = presentation.slides.add_slide(bullet_layout)
    deletion_slide.shapes.title.text = "Deletion Process"
    deletion_body = deletion_slide.placeholders[1].text_frame
    deletion_body.text = "An automated nightly job purges records past their retention window."
    deletion_body.add_paragraph().text = "No manual review step -- the schedule is enforced consistently."

    presentation.save(str(OUT_DIR / "data_retention_policy.pptx"))


# ---------------------------------------------------------------------------
# Markdown -- an informal internal quick-reference note.
# ---------------------------------------------------------------------------

MARKDOWN_CONTENT = f"""# Data Retention Quick Reference

If you just need the numbers without reading the full policy document:

- {_fact_sentence("financial")}
- {_fact_sentence("employee")}
- {_fact_sentence("support")}

See the full Data Retention Policy document for the formal deletion
process. This note is a personal quick-reference, not the source of
truth.
"""


def build_markdown() -> None:
    (OUT_DIR / "data_retention_policy.md").write_text(MARKDOWN_CONTENT, encoding="utf-8")


# ---------------------------------------------------------------------------
# Code -- a small script that would plausibly implement the same policy.
# ---------------------------------------------------------------------------

CODE_CONTENT = '''"""
Nightly data retention purge job.

Implements the schedule described in the company's Data Retention Policy:
financial records are kept 7 years, employee records 3 years after
termination, and customer support tickets 2 years, after which each
category is deleted by this job rather than by manual review.
"""

FINANCIAL_RECORDS_RETENTION_YEARS = 7
EMPLOYEE_RECORDS_RETENTION_YEARS_POST_TERMINATION = 3
SUPPORT_TICKET_RETENTION_YEARS = 2


def purge_expired_records(records, retention_years, reference_date):
    """Returns the subset of `records` older than `retention_years` relative
    to `reference_date`. A real deployment would delete these; this stub
    only selects them, so the policy stays auditable before anything is
    actually removed."""
    cutoff_year = reference_date.year - retention_years
    return [r for r in records if r.get("year", cutoff_year) < cutoff_year]
'''


def build_code() -> None:
    (OUT_DIR / "data_retention_policy.py").write_text(CODE_CONTENT, encoding="utf-8")


# ---------------------------------------------------------------------------
# Image -- a rendered "photograph of a printed policy excerpt" for the OCR
# extractor. Synthetic (not real handwriting -- deterministic handwriting
# generation isn't practical), but genuinely run through the real
# ImageOcrExtractor/Tesseract path at ingestion time, same as a scanned page
# would be (see docs/adr/0012-multi-format-extraction.md).
# ---------------------------------------------------------------------------

IMAGE_LINES = [
    "Data Retention Policy",
    "Financial records: retained seven years.",
    "Employee records: retained three years",
    "after termination.",
]

_FONT_CANDIDATES = [
    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
]


def _load_font(size: int) -> ImageFont.ImageFont:
    for path in _FONT_CANDIDATES:
        if Path(path).exists():
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def build_image() -> None:
    width, height = 900, 320
    image = Image.new("RGB", (width, height), color="white")
    draw = ImageDraw.Draw(image)
    font = _load_font(32)
    y = 30
    for line in IMAGE_LINES:
        draw.text((40, y), line, fill="black", font=font)
        y += 60
    image.save(OUT_DIR / "data_retention_policy.png")


if __name__ == "__main__":
    build_docx()
    print("wrote data_retention_policy.docx")
    build_pptx()
    print("wrote data_retention_policy.pptx")
    build_markdown()
    print("wrote data_retention_policy.md")
    build_code()
    print("wrote data_retention_policy.py")
    build_image()
    print("wrote data_retention_policy.png")
