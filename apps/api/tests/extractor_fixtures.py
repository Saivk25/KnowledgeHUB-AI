"""
Test helpers for building real DOCX/PPTX/image fixtures (Milestone 5).

Same philosophy as tests/pdf_helpers.py: build real files with the actual
library (python-docx, python-pptx, Pillow) so extraction tests hit the real
parsing/OCR code path rather than a mock. Each of these libraries is a
Milestone 5 dependency (see requirements.txt), not a Milestone 1 one, so
callers should still be fine importing this module directly in Milestone
5-only test files (no importorskip needed at the top of this file itself,
since every test file that imports it already requires these packages to
exist for the extractor code under test to even import).
"""

from __future__ import annotations


def make_sample_docx(path: str, paragraphs: list[str]) -> None:
    import docx

    document = docx.Document()
    for text in paragraphs:
        document.add_paragraph(text)
    document.save(path)


def make_sample_pptx(path: str, slide_texts: list[str]) -> None:
    import pptx

    presentation = pptx.Presentation()
    layout = presentation.slide_layouts[1]  # "Title and Content"
    for text in slide_texts:
        slide = presentation.slides.add_slide(layout)
        slide.shapes.title.text = text
    presentation.save(path)


def make_sample_image(path: str, text: str) -> None:
    """A clean, printed-text PNG -- deliberately not handwriting. Printed
    text OCRs reliably and deterministically, which is what a CI-safe test
    needs; handwriting recognition quality is explicitly out of scope for
    this milestone (see ADR-0012 / ADR-0006)."""

    from PIL import Image, ImageDraw, ImageFont

    image = Image.new("RGB", (800, 200), color="white")
    draw = ImageDraw.Draw(image)
    try:
        font = ImageFont.load_default(size=48)
    except TypeError:  # Pillow < 10.1 doesn't support the size kwarg
        font = ImageFont.load_default()
    draw.text((20, 60), text, fill="black", font=font)
    image.save(path)
