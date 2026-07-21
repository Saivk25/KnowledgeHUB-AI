"""
Extraction: an Extractor registry, one implementation per source format.

Milestone 3 shipped one function, `extract_text(pdf_path) -> pages`, hardcoded
to PyMuPDF. Milestone 5 (Multi-Format Ingestion) generalizes this into the
same kind of pluggable-provider registry the codebase already uses for
`EmbeddingProvider` (services/embeddings.py) and `VectorRepository`
(services/vector_repo.py): one `Extractor` interface, one concrete class per
format, resolved from a filename's extension. Adding a future format (e.g.
audio transcription, per Vision v2 Phase 2 capture types) means adding one
new class and one registry entry -- nothing in ingestion_service.py or
api/v1/routes/documents.py needs to change to recognize it, beyond adding the
extension to the upload allowlist.

Normalized shape: every extractor returns the same `ExtractionResult`,
built from a list of `ExtractedUnit` -- a "page-equivalent" (text +
sequential position + confidence). This is deliberately a smaller, MVP-scoped
version of Vision v2's `ContentBlock` idea (Section 4): exactly what M5 needs
(text + position + confidence), not the fuller shape a later capture
milestone might add (timestamps, image references, etc.).

Naming decision (see docs/adr/0012-multi-format-extraction.md): the unit
position is still called `page_number` everywhere downstream (ResourcePage,
ResourceChunk, the citation/retrieval modules, the frontend Source Viewer) --
not renamed to something generic like `unit_number`. For PDF it's a real
page; for PPTX it's a slide number; for every single-unit format (DOCX,
TXT/Markdown, code files, a fetched YouTube transcript) it's always `1`.
Extending the meaning while keeping the name matches the precedent Milestone
4 already set for `VectorPoint.document_id`.

Confidence: 1.0 for every deterministic parser (DOCX/PPTX/TXT/MD/code -- there
is no ambiguity in reading a zip/XML/plain-text file). Only image OCR
produces a confidence below 1.0, and it is always the OCR engine's own
reported score, never an invented heuristic (DRR Section 9).
"""

from __future__ import annotations

from abc import ABC, abstractmethod
from pathlib import Path

from app.core.config import get_settings

settings = get_settings()

MIN_CHARS_PER_RESOURCE = 20


class ExtractionError(Exception):
    """Raised by an Extractor when a file can't be parsed at all (corrupt
    zip/XML, unreadable image, etc.). Caught in ingestion_service.py and
    turned into a FAILED status with `code`/`message`, the same pattern
    already established for SCANNED_PDF_UNSUPPORTED."""

    def __init__(self, code: str, message: str):
        self.code = code
        self.message = message
        super().__init__(message)


class ExtractedUnit:
    def __init__(self, unit_number: int, text: str, confidence: float = 1.0):
        self.unit_number = unit_number
        self.text = text
        self.confidence = confidence


class ExtractionResult:
    def __init__(self, units: list[ExtractedUnit]):
        self.units = units

    @property
    def pages(self) -> list[tuple[int, str]]:
        """Backward-compatible view: chunking.py and ingestion_service.py's
        callers keep working unchanged against a (page_number, text) shape,
        regardless of which extractor produced the units."""
        return [(u.unit_number, u.text) for u in self.units]

    @property
    def page_count(self) -> int:
        return len(self.units)

    @property
    def total_chars(self) -> int:
        return sum(len(u.text) for u in self.units)

    @property
    def looks_scanned(self) -> bool:
        return self.total_chars < MIN_CHARS_PER_RESOURCE

    @property
    def confidence(self) -> float:
        """Aggregate extraction confidence stored on Resource.extraction_confidence.
        A simple average across units -- every non-OCR extractor's units are
        all 1.0 already, so this only does real work for OCR."""
        if not self.units:
            return 1.0
        return sum(u.confidence for u in self.units) / len(self.units)


class Extractor(ABC):
    #: File extensions (lowercase, with leading dot) this extractor handles.
    extensions: frozenset[str] = frozenset()

    @abstractmethod
    def extract(self, path: str) -> ExtractionResult: ...


class PdfExtractor(Extractor):
    """PyMuPDF native text layer only -- unchanged from Milestone 3
    (see ADR-0006). Scanned/image-only PDFs are still reported via
    `looks_scanned`, not silently indexed as empty content."""

    extensions = frozenset({".pdf"})

    def extract(self, path: str) -> ExtractionResult:
        import fitz  # PyMuPDF

        try:
            doc = fitz.open(path)
        except Exception as exc:  # noqa: BLE001
            raise ExtractionError("UNREADABLE_PDF", "This PDF could not be opened.") from exc
        units: list[ExtractedUnit] = []
        try:
            for index in range(doc.page_count):
                page = doc.load_page(index)
                text = page.get_text("text") or ""
                units.append(ExtractedUnit(unit_number=index + 1, text=text.strip()))
        finally:
            doc.close()
        return ExtractionResult(units)


class DocxExtractor(Extractor):
    """python-docx has no reliable notion of a rendered "page" (page breaks
    depend on the reader's layout engine, not the file format), so a DOCX
    resource is always exactly one unit: every paragraph's text, joined with
    newlines, in document order."""

    extensions = frozenset({".docx"})

    def extract(self, path: str) -> ExtractionResult:
        import docx
        from docx.opc.exceptions import PackageNotFoundError

        try:
            document = docx.Document(path)
        except (PackageNotFoundError, Exception) as exc:  # noqa: BLE001
            raise ExtractionError("UNREADABLE_DOCX", "This Word document could not be opened.") from exc
        text = "\n".join(p.text for p in document.paragraphs)
        return ExtractionResult([ExtractedUnit(unit_number=1, text=text.strip())])


class PptxExtractor(Extractor):
    """One unit per slide (text from every shape/text frame on that slide),
    the natural PPTX analogue of a PDF page -- slide number is what a Source
    Viewer would cite, same shape as PDF page citations."""

    extensions = frozenset({".pptx"})

    def extract(self, path: str) -> ExtractionResult:
        import pptx
        from pptx.exc import PackageNotFoundError

        try:
            presentation = pptx.Presentation(path)
        except (PackageNotFoundError, Exception) as exc:  # noqa: BLE001
            raise ExtractionError("UNREADABLE_PPTX", "This PowerPoint file could not be opened.") from exc
        units: list[ExtractedUnit] = []
        for index, slide in enumerate(presentation.slides):
            fragments = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    fragments.append(shape.text_frame.text)
            units.append(ExtractedUnit(unit_number=index + 1, text="\n".join(fragments).strip()))
        return ExtractionResult(units)


class TextExtractor(Extractor):
    """Plain text and Markdown: read the whole file as one unit. No
    heading-aware sectioning (that is an M8-era chunking-strategy question,
    per Vision v2 Section 2/Architecture Section 6 -- out of this milestone's
    scope, see the Milestone 5 design note). Also the extractor used for a
    fetched YouTube transcript, which is cached to disk as a plain .txt file
    through the existing storage service and then flows through this exact
    path -- no separate transcript-aware extractor needed."""

    extensions = frozenset({".txt", ".md"})

    def extract(self, path: str) -> ExtractionResult:
        try:
            text = Path(path).read_text(encoding="utf-8", errors="replace")
        except OSError as exc:
            raise ExtractionError("UNREADABLE_FILE", "This file could not be read.") from exc
        return ExtractionResult([ExtractedUnit(unit_number=1, text=text.strip())])


class CodeExtractor(Extractor):
    """Source code files: read as plain text, one unit per file (no
    function/class-aware chunking -- tree-sitter-based semantic chunking is
    explicitly out of scope for this milestone, per the approved design)."""

    extensions = frozenset(
        {
            ".py",
            ".js",
            ".jsx",
            ".ts",
            ".tsx",
            ".java",
            ".c",
            ".h",
            ".cpp",
            ".hpp",
            ".go",
            ".rs",
            ".rb",
            ".php",
            ".cs",
            ".kt",
            ".sql",
            ".sh",
        }
    )

    def extract(self, path: str) -> ExtractionResult:
        try:
            text = Path(path).read_text(encoding="utf-8", errors="replace")
        except OSError as exc:
            raise ExtractionError("UNREADABLE_FILE", "This file could not be read.") from exc
        return ExtractionResult([ExtractedUnit(unit_number=1, text=text.strip())])


class ImageOcrExtractor(Extractor):
    """OCR via pytesseract + the Tesseract engine (see
    docs/adr/0012-multi-format-extraction.md, superseding ADR-0006's
    "no OCR in MVP" for this one format). Confidence is Tesseract's own
    mean per-word confidence (0-100), normalized to 0-1 -- never an invented
    number (DRR Section 9). Handwriting recognition is not solved by this;
    the honest, surfaced confidence score is the product's answer to that
    limitation (Architecture doc risk #2), not a claim of accuracy."""

    extensions = frozenset({".png", ".jpg", ".jpeg"})

    def extract(self, path: str) -> ExtractionResult:
        import pytesseract
        from PIL import Image, UnidentifiedImageError

        if settings.TESSERACT_CMD:
            pytesseract.pytesseract.tesseract_cmd = settings.TESSERACT_CMD

        try:
            image = Image.open(path)
            image.load()
        except (UnidentifiedImageError, OSError) as exc:
            raise ExtractionError("UNREADABLE_IMAGE", "This image could not be opened.") from exc

        try:
            data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
        except pytesseract.TesseractNotFoundError as exc:
            raise ExtractionError(
                "OCR_ENGINE_UNAVAILABLE",
                "The OCR engine is not installed on this server.",
            ) from exc

        words = [w for w in data.get("text", []) if w and w.strip()]
        confidences = [
            int(c)
            for w, c in zip(data.get("text", []), data.get("conf", []), strict=False)
            if w and w.strip() and str(c).lstrip("-").isdigit() and int(c) >= 0
        ]
        text = " ".join(words)
        confidence = (sum(confidences) / len(confidences) / 100.0) if confidences else 0.0
        return ExtractionResult([ExtractedUnit(unit_number=1, text=text.strip(), confidence=confidence)])


# Registry: order doesn't matter, extensions are disjoint across extractors.
_EXTRACTORS: list[Extractor] = [
    PdfExtractor(),
    DocxExtractor(),
    PptxExtractor(),
    TextExtractor(),
    CodeExtractor(),
    ImageOcrExtractor(),
]

#: Every extension any extractor supports -- used by the upload route for
#: allowlist validation and by the frontend-facing error message.
SUPPORTED_EXTENSIONS: frozenset[str] = frozenset().union(*(e.extensions for e in _EXTRACTORS))

_MIME_TYPES: dict[str, str] = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
}


def _extension_of(filename: str) -> str:
    return Path(filename).suffix.lower()


def get_extractor_for(filename: str) -> Extractor | None:
    ext = _extension_of(filename)
    for extractor in _EXTRACTORS:
        if ext in extractor.extensions:
            return extractor
    return None


def mime_type_for(filename: str) -> str:
    return _MIME_TYPES.get(_extension_of(filename), "text/plain")


def is_supported_filename(filename: str) -> bool:
    return _extension_of(filename) in SUPPORTED_EXTENSIONS


# -- Backward-compatible function, kept for anything importing the old
# Milestone 3 name directly (nothing in this codebase does after this
# milestone's edits, but this keeps the module's public surface additive
# rather than breaking). Delegates to PdfExtractor.
def extract_text(pdf_path: str) -> ExtractionResult:
    return PdfExtractor().extract(pdf_path)
